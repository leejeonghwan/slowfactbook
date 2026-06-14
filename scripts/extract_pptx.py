#!/usr/bin/env python3
"""
PowerPoint (.pptx) -> structured JSON extractor (Slow Factbook).

Usage:
  python3 extract_pptx.py input.pptx output.json
Requires: python-pptx  (pip install python-pptx)

Alternative to extract_keynote.py. PPTX stores chart data as an embedded
spreadsheet, so python-pptx reads series/categories/values directly — and
pie/doughnut charts (which the Keynote parser chokes on) extract cleanly.
CAVEAT: only works if Keynote exported charts as NATIVE charts. Any chart
that exported as a flat image has NO data and is reported under "_image_only".
"""
import sys, os, json, re

CT_MAP = {
    "LINE": "line", "LINE_MARKERS": "line", "LINE_STACKED": "line",
    "COLUMN_CLUSTERED": "column", "COLUMN_STACKED": "stacked_bar",
    "COLUMN_STACKED_100": "stacked_bar",
    "BAR_CLUSTERED": "bar", "BAR_STACKED": "stacked_bar_h",
    "BAR_STACKED_100": "stacked_bar_h",
    "PIE": "pie", "PIE_EXPLODED": "pie",
    "DOUGHNUT": "pie", "DOUGHNUT_EXPLODED": "pie",
    "XY_SCATTER": "scatter", "AREA": "area", "AREA_STACKED": "area",
}

UNIT_RE = re.compile(r"단위|출처|통계청|국가데이터처|데이터처|노동부|고용노동부|보건복지부|"
                     r"기획재정부|행정안전부|한국은행|OECD|World Bank|World|기준|조사|"
                     r"연구원|연구소|재단|협회|학회|위원회|진흥원|개발원|평가원|"
                     r"은행|보험|증권|일보|신문|방송|IMF|UN|WHO|ILO")
URL_RE = re.compile(r"https?://")
# a bare value / axis label like "1500만 명.", "500만 명.", "2024", "12.3%" — NOT a title
VALUE_RE = re.compile(r"^[\d][\d,.\s]*\s*(만|억|천|조)?\s*"
                      r"(명|원|개|건|가구|채|세대|시간|년|개월|배|달러|톤|%|p)?\s*\.?$")
# short chart annotations (region/series labels), often placed above the real
# title — e.g. "60세 이전.", "60세 이후.", "65세 이상.", "남성.", "여성." — NOT a title
LABEL_RE = re.compile(r"(이전|이후|이상|미만|초과|이내|이하)\s*\.?$"
                      r"|^\d+\s*세\s*\.?$"
                      r"|^(남성|여성|남자|여자|남|여)\s*\.?$"
                      r"|^\S{1,6}\s+-?\d[\d,.]*\s*%?\s*\.?$")  # data callout like "한국 4.3"

def title_placeholder(slide):
    """If the slide uses a real Title placeholder, return its text — the most
    reliable signal. Slides authored to this convention extract 100% correctly."""
    for sh in slide.shapes:
        try:
            if sh.is_placeholder and sh.has_text_frame:
                t = sh.placeholder_format.type
                if t is not None and "TITLE" in str(t):
                    txt = " ".join(sh.text_frame.text.split())
                    if txt:
                        return txt
        except Exception:
            pass
    return None

def pick_title_source(slide):
    """Title priority: (1) Title placeholder, else (2) top-most text box after
    dropping source lines, bare value/axis labels, and short chart annotations.
    Source = line matching UNIT_RE."""
    cand = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            top = sh.top if sh.top is not None else 10**9
            cand.append((top, " ".join(sh.text_frame.text.split())))
    cand.sort(key=lambda x: x[0])
    urls = [t for _, t in cand if URL_RE.search(t)]
    url = urls[0] if urls else ""
    srcs = [t for _, t in cand if UNIT_RE.search(t)]
    source = srcs[0] if srcs else ""
    ph = title_placeholder(slide)
    if ph:
        return ph, source, url
    base = [t for _, t in cand
            if not UNIT_RE.search(t) and not URL_RE.search(t) and not VALUE_RE.match(t)]
    strong = [t for t in base if not LABEL_RE.search(t)]   # drop chart annotations
    # never fall back to a URL/source as the title — leave it empty so the chart's
    # own title (e.g. doughnut center) can take over.
    title = (strong or base or [""])[0]
    return title, source, url

def review_title(slide, chosen):
    """Flag slides whose title is uncertain, so they can be reviewed/overridden
    instead of eyeballing the whole deck. None = confident."""
    if title_placeholder(slide):
        return None
    cand = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            cand.append(" ".join(sh.text_frame.text.split()))
    base = [t for t in cand if not UNIT_RE.search(t) and not URL_RE.search(t) and not VALUE_RE.match(t)]
    strong = [t for t in base if not LABEL_RE.search(t)]
    if not (chosen or "").strip():
        return {"reason": "no_title", "candidates": cand[:6]}
    if len(strong) >= 2:                      # genuinely ambiguous
        return {"reason": "ambiguous", "chosen": chosen, "candidates": strong[:6]}
    return None

def _base_title(t):
    return re.sub(r"\s*\(\d+\)\s*$", "", t).strip()

def _richness(it):
    return sum(1 for s in it.get("series", []) for v in s if v is not None)

def dedup_builds(items):
    """Keynote build animations export as consecutive duplicate slides.
    Collapse each run of consecutive same-title charts to its richest instance."""
    out, run = [], []
    def flush():
        if run:
            out.append(max(run, key=_richness))
    last = None
    for it in items:
        if not it.get("vizType"):
            flush(); run.clear(); out.append(it); last = None; continue
        bt = _base_title(it["title"])
        if bt == last:
            run.append(it)
        else:
            flush(); run.clear(); run.append(it); last = bt
    flush()
    return out

MAX_POINTS = 1500  # downsample very large series (e.g. daily decades-long data)

def _cache_points(container):
    """Return {idx: text} from a c:cat/c:val cache, parsed directly from XML (fast)."""
    from pptx.oxml.ns import qn
    d = {}
    if container is None:
        return d
    for pt in container.iter(qn('c:pt')):
        try:
            idx = int(pt.get('idx', 0))
        except (TypeError, ValueError):
            idx = 0
        v = pt.find(qn('c:v'))
        d[idx] = v.text if (v is not None and v.text is not None) else None
    return d

def _to_num(s):
    if s is None:
        return None
    try:
        return float(s)
    except ValueError:
        return None

def chart_title(chart):
    """The chart's own title (often shown in a doughnut's center). Many slides
    put the real title here instead of in a text box."""
    try:
        if chart.has_title:
            return " ".join(chart.chart_title.text_frame.text.split())
    except Exception:
        pass
    return ""

def is_weak_title(title, source):
    """True when the slide text-box title isn't a real title (missing, equals the
    source line, or is just a unit/value/annotation) — then use the chart title."""
    t = (title or "").strip()
    return (not t) or t == (source or "").strip() or bool(UNIT_RE.search(t)) \
        or bool(VALUE_RE.match(t)) or bool(LABEL_RE.search(t))

def parse_chart_xml(chart):
    """Fast chart extraction via raw XML. Avoids python-pptx's slow
    categories/series.values, which hang on charts with tens of thousands
    of points. Returns (seriesNames, labels, series_values)."""
    from pptx.oxml.ns import qn
    cs = chart._chartSpace
    sers = list(cs.iter(qn('c:ser')))
    names, series, label_map, max_idx = [], [], {}, 0
    for ser in sers:
        tx = ser.find(qn('c:tx'))
        nm = None
        if tx is not None:
            vv = tx.find('.//' + qn('c:v'))
            nm = vv.text if vv is not None else None
        names.append(nm)
        vals = _cache_points(ser.find(qn('c:val')))
        cats = _cache_points(ser.find(qn('c:cat')))
        if cats and len(cats) > len(label_map):
            label_map = cats
        n = (max(vals) if vals else -1)
        max_idx = max(max_idx, n, (max(cats) if cats else -1))
        series.append(vals)
    L = max_idx + 1
    labels = [label_map.get(i, "") for i in range(L)]
    out_series = [[_to_num(sv.get(i)) for i in range(L)] for sv in series]
    # downsample very large series for web rendering
    if L > MAX_POINTS:
        step = (L + MAX_POINTS - 1) // MAX_POINTS
        keep = list(range(0, L, step))
        if keep[-1] != L - 1:
            keep.append(L - 1)
        labels = [labels[i] for i in keep]
        out_series = [[s[i] for i in keep] for s in out_series]
    return names, labels, out_series

def series_axes_kinds(chart):
    """For combo/dual-axis charts: per series (in the same order as parse_chart_xml),
    return its kind ('bar'/'line'/'area') and value-axis index (0=left, 1=right)."""
    from pptx.oxml.ns import qn
    cs = chart._chartSpace
    valax_ids = []
    for ax in cs.findall(".//" + qn("c:valAx")):
        a = ax.find(qn("c:axId"))
        if a is not None:
            valax_ids.append(a.get("val"))
    kinds, axes = [], []
    for ser in cs.iter(qn("c:ser")):
        grp = ser.getparent()
        gtag = grp.tag.split("}")[1].lower() if grp is not None else ""
        kind = "bar" if "bar" in gtag else ("area" if "area" in gtag else "line")
        axids = [a.get("val") for a in grp.findall(qn("c:axId"))] if grp is not None else []
        val_axid = next((a for a in axids if a in valax_ids), None)
        axes.append(valax_ids.index(val_axid) if val_axid in valax_ids else 0)
        kinds.append(kind)
    return kinds, axes

def extract(inp, outp, start=1, end=None, dedup=True):
    from pptx import Presentation
    prs = Presentation(inp)
    all_slides = list(prs.slides)
    if end is None:
        end = len(all_slides)
    items, image_only, title_review = [], [], []
    current_section = ""   # section dividers use the 'Title' layout
    for si in range(start, end + 1):
        slide = all_slides[si - 1]
        layout = slide.slide_layout.name
        title, source, surl = pick_title_source(slide)
        charts = [sh for sh in slide.shapes if sh.has_chart]
        pics = [sh for sh in slide.shapes if sh.shape_type == 13]  # PICTURE
        if not charts:
            if layout == "Title" and title.strip():
                current_section = title.strip().rstrip(".")
            if pics:
                image_only.append({"slide": si, "title": title})
            items.append({"slide": f"slide-{si}", "title": title, "source": source,
                          "sourceUrl": surl, "vizType": None, "layout": layout,
                          "category": current_section, "note": "no native chart"})
            continue
        weak = is_weak_title(title, source)
        any_chart_title = any(chart_title(sh.chart) for sh in charts)
        rv = review_title(slide, title)
        if rv and not (weak and any_chart_title):   # chart title rescues it
            rv["slide"] = f"slide-{si}"
            title_review.append(rv)
        for ci, sh in enumerate(charts):
            ch = sh.chart
            ct = str(ch.chart_type).split()[0] if ch.chart_type else ""
            viz = CT_MAP.get(ct, ct.lower())
            names, labels, series = parse_chart_xml(ch)
            kinds, axes = series_axes_kinds(ch)
            if (len(set(kinds)) > 1) or (axes and max(axes) > 0):
                viz = "combo"   # mixed bar+line and/or dual value axis
            cht = chart_title(ch)
            # RULE: pie/doughnut -> always use the chart's own (center) title first.
            # Other types -> use the chart title only when the slide title is weak.
            if cht and (viz == "pie" or weak):
                full_title = cht
            else:
                # multiple charts share the slide title; distinguish by series name
                # ("...영업이익 — 반도체"), else a numeric suffix.
                suffix = ""
                if len(charts) > 1:
                    distinct = [n.strip().rstrip(".") for n in names if n and n.strip()]
                    suffix = " — " + distinct[0] if len(distinct) == 1 else f" ({ci+1})"
                full_title = title + suffix
            item = {
                "slide": f"slide-{si}",
                "title": full_title,
                "source": source,
                "sourceUrl": surl,
                "category": current_section,
                "vizType": viz,
                "chartType_raw": ct,
                "seriesNames": names,
                "labels": [str(c).strip().rstrip(".") for c in labels],
                "series": series,
            }
            if viz == "combo":
                item["seriesKinds"] = kinds
                item["seriesAxes"] = axes
            items.append(item)
    raw_count = sum(1 for i in items if i.get("vizType"))
    if dedup:
        items = dedup_builds(items)
    out = {"_source_file": os.path.basename(inp),
           "_slides": len(all_slides), "_range": [start, end],
           "_charts": sum(1 for i in items if i.get("vizType")),
           "_charts_before_dedup": raw_count,
           "_image_only": image_only,
           "_title_review": title_review, "items": items}
    json.dump(out, open(outp, "w", encoding="utf-8"), ensure_ascii=False)
    print(f"range={start}-{end} charts={out['_charts']} image_only={len(image_only)} -> {outp}")

if __name__ == "__main__":
    a = sys.argv
    if len(a) >= 5:   # input out start end  (chunk mode, no per-chunk dedup)
        extract(a[1], a[2], int(a[3]), int(a[4]), dedup=False)
    else:
        extract(a[1], a[2])
